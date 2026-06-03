from pathlib import Path
from .config import DEFAULT_MODEL, parse_provider, get_default_model
from .providers import create_message
import chromadb
from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction

CHUNK_SIZE = 150  # words per chunk
CHUNK_OVERLAP = 20


def _load_file(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        import pypdf
        reader = pypdf.PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    elif suffix == ".docx":
        import docx
        doc = docx.Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    elif suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(path))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text)
        return "\n".join(lines)
    elif suffix == ".ipynb":
        import nbformat
        nb = nbformat.read(str(path), as_version=4)
        lines = []
        for cell in nb.cells:
            if cell.cell_type in ("markdown", "code") and cell.source.strip():
                lines.append(cell.source)
        return "\n\n".join(lines)
    else:
        return path.read_text(encoding="utf-8")


def _chunk(text: str) -> list[str]:
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunks.append(" ".join(words[i:i + CHUNK_SIZE]))
        i += CHUNK_SIZE - CHUNK_OVERLAP
    return [c for c in chunks if c.strip()]


DEFAULT_SYSTEM_PROMPT = (
    "You are a senior data scientist answering questions based on analysis reports. "
    "Answer only from the provided context. Be specific and cite findings where relevant. "
    "If the context does not contain enough information to answer, say so clearly."
)


def ask(question: str, source: str = None, text: str = None, provider: str = None, model: str = None, system_prompt: str = None) -> str:
    """
    Ask a question across a collection of analysis documents or raw text.

    Args:
        question:      The question to answer.
        source:        Path to a folder containing .txt, .md, .pdf, .docx, .pptx, or .ipynb files.
        text:          A raw text string to search instead of a folder.
        provider:      Optional. The AI provider to use ("anthropic", "openai", "gemini").
                       If not specified, defaults to "anthropic" or infers from model.
        model:         Optional. The specific model to use. If not specified, uses the provider's default.
        system_prompt: Optional. A custom system prompt to override the default answering persona.

    Returns:
        An answer grounded in the provided documents.
    """
    if not source and not text:
        raise ValueError("Provide either 'source' (folder path) or 'text'.")

    # Parse provider and determine model
    provider_enum = parse_provider(provider, model)
    if model is None:
        model = get_default_model(provider_enum)

    # Collect chunks
    chunks = []

    if text:
        chunks.extend(_chunk(text))

    if source:
        folder = Path(source).expanduser().resolve()
        supported = {".txt", ".md", ".pdf", ".docx", ".pptx", ".ipynb"}
        for file in sorted(folder.iterdir()):
            if file.suffix.lower() in supported:
                content = _load_file(file)
                chunks.extend(_chunk(content))

    if not chunks:
        raise ValueError("No content found. Check your source folder or text input.")

    # Embed and store in ChromaDB
    embedding_fn = SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")
    client = chromadb.Client()
    collection = client.get_or_create_collection(
        name="bridgekit_ask",
        embedding_function=embedding_fn
    )
    collection.add(
        documents=chunks,
        ids=[f"chunk_{i}" for i in range(len(chunks))]
    )

    # Retrieve most relevant chunks
    results = collection.query(query_texts=[question], n_results=min(8, len(chunks)))
    context = "\n\n".join(results["documents"][0])

    # Generate answer with specified provider
    user_message = f"Context from analysis reports:\n\n{context}\n\nQuestion: {question}"

    return create_message(
        provider=provider_enum,
        system_prompt=system_prompt or DEFAULT_SYSTEM_PROMPT,
        user_message=user_message,
        model=model,
        max_tokens=1024
    )
